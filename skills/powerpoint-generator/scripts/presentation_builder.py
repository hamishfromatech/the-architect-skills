"""
Presentation Builder - Generate PowerPoint presentations from structured data.

This module provides utilities for creating presentations with slides,
text, images, tables, and charts.
"""

from dataclasses import dataclass
from typing import List, Optional, Dict, Any, Union
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


@dataclass
class SlideContent:
    """Content for a single slide."""
    title: str
    subtitle: Optional[str] = None
    content: List[str] = None
    images: List[Dict[str, Any]] = None
    table_data: Optional[Dict[str, Any]] = None
    chart_data: Optional[Dict[str, Any]] = None

    def __post_init__(self):
        if self.content is None:
            self.content = []
        if self.images is None:
            self.images = []


@dataclass
class PresentationConfig:
    """Configuration for presentation generation."""
    title: str
    slides: List[SlideContent]
    author: Optional[str] = None
    company: Optional[str] = None
    aspect_ratio: str = "16:9"  # or "4:3"


class PresentationBuilder:
    """Build PowerPoint presentations from structured data."""

    # Standard slide layouts
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TITLE_ONLY = 5
    LAYOUT_BLANK = 6

    def __init__(self, config: PresentationConfig):
        """Initialize the presentation builder.

        Args:
            config: Presentation configuration with slides data
        """
        self.config = config
        self.prs = Presentation()

        # Set aspect ratio if 4:3
        if config.aspect_ratio == "4:3":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)

    def build(self) -> Presentation:
        """Build the complete presentation.

        Returns:
            Presentation object ready to save
        """
        for slide_content in self.config.slides:
            self._add_slide(slide_content)
        return self.prs

    def _add_slide(self, content: SlideContent) -> None:
        """Add a slide with the given content.

        Args:
            content: Slide content specification
        """
        # Determine layout
        if content.subtitle and not content.content:
            layout_idx = self.LAYOUT_TITLE
        elif content.table_data or content.chart_data:
            layout_idx = self.LAYOUT_TITLE_ONLY
        else:
            layout_idx = self.LAYOUT_TITLE_CONTENT

        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.title

        # Set subtitle (for title slide)
        if content.subtitle and layout_idx == self.LAYOUT_TITLE:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content.subtitle

        # Add bullet content
        if content.content and layout_idx == self.LAYOUT_TITLE_CONTENT:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            for i, bullet in enumerate(content.content):
                if i == 0:
                    tf.text = bullet
                else:
                    tf.add_paragraph().text = bullet

        # Add images
        for img in content.images:
            slide.shapes.add_picture(
                img['path'],
                Inches(img.get('left', 1)),
                Inches(img.get('top', 1)),
                width=Inches(img['width']) if 'width' in img else None,
                height=Inches(img['height']) if 'height' in img else None
            )

        # Add table
        if content.table_data:
            self._add_table(slide, content.table_data)

        # Add chart
        if content.chart_data:
            self._add_chart(slide, content.chart_data)

    def _add_table(self, slide, table_data: Dict[str, Any]) -> None:
        """Add a table to a slide.

        Args:
            slide: Slide object
            table_data: Table configuration with headers and rows
        """
        rows = len(table_data['rows']) + 1  # +1 for header
        cols = len(table_data['headers'])

        x = Inches(table_data.get('left', 2))
        y = Inches(table_data.get('top', 2))
        cx = Inches(table_data.get('width', 6))
        cy = Inches(table_data.get('height', 1.5))

        shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)
        table = shape.table

        # Set headers
        for col_idx, header in enumerate(table_data['headers']):
            table.cell(0, col_idx).text = header

        # Set data rows
        for row_idx, row in enumerate(table_data['rows']):
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)

    def _add_chart(self, slide, chart_data: Dict[str, Any]) -> None:
        """Add a chart to a slide.

        Args:
            slide: Slide object
            chart_data: Chart configuration with categories and series
        """
        data = CategoryChartData()
        data.categories = chart_data['categories']

        for series_name, values in chart_data['series'].items():
            data.add_series(series_name, values)

        chart_type_str = chart_data.get('type', 'column')
        chart_type = self._get_chart_type(chart_type_str)

        x = Inches(chart_data.get('left', 2))
        y = Inches(chart_data.get('top', 2))
        cx = Inches(chart_data.get('width', 6))
        cy = Inches(chart_data.get('height', 4.5))

        slide.shapes.add_chart(chart_type, x, y, cx, cy, data)

    def _get_chart_type(self, chart_type: str) -> XL_CHART_TYPE:
        """Get chart type enum from string.

        Args:
            chart_type: String representation of chart type

        Returns:
            XL_CHART_TYPE enum value
        """
        chart_map = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'area': XL_CHART_TYPE.AREA,
        }
        return chart_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    def save(self, filepath: Union[str, Path]) -> None:
        """Save the presentation to a file.

        Args:
            filepath: Output file path
        """
        self.prs.save(filepath)


def create_presentation(
    title: str,
    slides: List[Dict[str, Any]],
    output_path: str,
    **kwargs
) -> str:
    """Convenience function to create a presentation.

    Args:
        title: Presentation title
        slides: List of slide specifications
        output_path: Where to save the file
        **kwargs: Additional configuration options

    Returns:
        Path to saved presentation
    """
    slide_contents = [
        SlideContent(**slide) if isinstance(slide, dict) else slide
        for slide in slides
    ]

    config = PresentationConfig(
        title=title,
        slides=slide_contents,
        **kwargs
    )

    builder = PresentationBuilder(config)
    builder.build()
    builder.save(output_path)

    return output_path


if __name__ == "__main__":
    # Example usage
    slides = [
        SlideContent(
            title="Quarterly Report",
            subtitle="Q4 2024 Results"
        ),
        SlideContent(
            title="Revenue Overview",
            content=[
                "Total revenue: $2.5M",
                "Year-over-year growth: 15%",
                "New customers: 150"
            ]
        ),
        SlideContent(
            title="Sales Data",
            chart_data={
                'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
                'series': {'Revenue': (1.8, 2.1, 2.3, 2.5)},
                'type': 'column'
            }
        )
    ]

    config = PresentationConfig(
        title="Business Report",
        slides=slides
    )

    builder = PresentationBuilder(config)
    builder.build()
    builder.save("example_report.pptx")
    print("Created example_report.pptx")